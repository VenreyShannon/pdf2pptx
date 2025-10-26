#!/usr/bin/env python3
"""
Images to PowerPoint Converter Module
"""

from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    Presentation = None
    Inches = None

try:
    from PIL import Image
    PILLOW_AVAILABLE = True
except ImportError:
    PILLOW_AVAILABLE = False
    Image = None


def convert_images_to_ppt(images_dir, output_ppt, slide_width=None, slide_height=None):
    """
    Convert images to PowerPoint presentation
    
    Args:
        images_dir: Directory containing the images
        output_ppt: Output PowerPoint file path
        slide_width: Width of slides in inches (default: 13.333 for 16:9)
        slide_height: Height of slides in inches (default: 7.5 for 16:9)
        
    Returns:
        Path to the generated PowerPoint file
    """
    if not PPTX_AVAILABLE:
        raise ImportError(
            "python-pptx library not found.\n"
            "Please install it using: pip install python-pptx"
        )
    
    if not PILLOW_AVAILABLE:
        raise ImportError(
            "Pillow library not found.\n"
            "Please install it using: pip install Pillow"
        )
    
    images_path = Path(images_dir)
    
    # Check if images directory exists
    if not images_path.exists():
        raise FileNotFoundError(f"Images directory '{images_dir}' not found!")
    
    # Get all PNG images sorted by name
    image_files = sorted(images_path.glob("*.png"))
    
    if not image_files:
        raise FileNotFoundError(f"No PNG images found in '{images_dir}'!")
    
    print(f"Found {len(image_files)} images")
    print(f"Creating PowerPoint presentation...")
    
    # Create presentation
    prs = Presentation()
    
    # Set slide size (16:9 aspect ratio)
    if slide_width and slide_height:
        prs.slide_width = Inches(slide_width)
        prs.slide_height = Inches(slide_height)
    else:
        # Default: 16:9 (standard widescreen)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    
    # Add each image as a slide
    for i, image_file in enumerate(image_files, start=1):
        print(f"  Adding slide {i}/{len(image_files)}: {image_file.name}")
        
        # Add blank slide
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Get image dimensions
        with Image.open(image_file) as img:
            img_width, img_height = img.size
        
        # Calculate image aspect ratio
        img_aspect = img_width / img_height
        slide_aspect = prs.slide_width / prs.slide_height
        
        # Calculate position and size to fit image in slide while maintaining aspect ratio
        if img_aspect > slide_aspect:
            # Image is wider than slide
            pic_width = prs.slide_width
            pic_height = int(prs.slide_width / img_aspect)
            left = 0
            top = (prs.slide_height - pic_height) // 2
        else:
            # Image is taller than slide
            pic_height = prs.slide_height
            pic_width = int(prs.slide_height * img_aspect)
            top = 0
            left = (prs.slide_width - pic_width) // 2
        
        # Add picture to slide
        slide.shapes.add_picture(
            str(image_file),
            left=left,
            top=top,
            width=pic_width,
            height=pic_height
        )
    
    # Save presentation
    output_path = Path(output_ppt)
    prs.save(str(output_path))
    
    print(f"✓ Successfully created PowerPoint presentation!")
    print(f"Total slides: {len(image_files)}")
    
    return str(output_path)

